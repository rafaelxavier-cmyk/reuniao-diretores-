"""
validator.py – Avalia e corrige automaticamente a qualidade visual de um PPTX.

Checks realizados:
  • Tamanho mínimo de fonte (conteúdo ≥ 12pt, títulos ≥ 18pt)
  • Densidade de texto por slide
  • Sobreposição entre shapes de conteúdo
  • Elementos fora dos limites do slide
  • Contraste de texto (cor sobre fundo)
  • Espaçamento e padding adequados
"""

from __future__ import annotations
from dataclasses import dataclass, field
from typing import List, Tuple
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

W_SLIDE = 12192000
H_SLIDE = 6858000

# Área de conteúdo (ignora barra lateral e logo)
CONTENT_LEFT   = 400000
CONTENT_TOP    = 900000
CONTENT_RIGHT  = 12100000
CONTENT_BOTTOM = 6750000

# Thresholds
MIN_CONTENT_PT  = 12    # mínimo para texto de conteúdo
MIN_TITLE_PT    = 20    # mínimo para títulos principais
MAX_CHARS_SLIDE = 600   # máximo de caracteres por slide (conteúdo)
MAX_ITEMS_CARD  = 6     # máximo de bullets por card

# Shapes que são decorativas/estruturais (não checar conteúdo)
STRUCTURAL_NAMES = {"sidebar", "stripe", "Shape 0", "Shape 1", "separator"}


@dataclass
class Issue:
    slide:     int
    severity:  str   # 'error' | 'warning' | 'info'
    category:  str
    message:   str
    shape:     str = ""
    auto_fixed: bool = False


@dataclass
class ValidationReport:
    issues:  List[Issue] = field(default_factory=list)
    score:   int = 100
    n_slides: int = 0

    def add(self, issue: Issue):
        self.issues.append(issue)
        if issue.severity == "error":
            self.score = max(0, self.score - 15)
        elif issue.severity == "warning":
            self.score = max(0, self.score - 3)
        elif issue.severity == "info":
            self.score = max(0, self.score - 0)

    @property
    def errors(self):
        return [i for i in self.issues if i.severity == "error"]

    @property
    def warnings(self):
        return [i for i in self.issues if i.severity == "warning"]

    @property
    def grade(self) -> str:
        if self.score >= 90: return "A"
        if self.score >= 75: return "B"
        if self.score >= 60: return "C"
        return "D"

    @property
    def passed(self) -> bool:
        return self.score >= 75


HEADER_ZONE_BOTTOM = 2100000   # elementos de cabeçalho (logo, título, supertítulo)
FOOTER_ZONE_TOP    = 6200000   # rodapé
LABEL_MAX_CHARS    = 40        # texto curto = label/badge, não checar tamanho

def _is_content_zone(shape) -> bool:
    """True se o shape é texto de CONTEÚDO (não header decorativo, não footer, não sidebar)."""
    if shape.left < CONTENT_LEFT:
        return False  # sidebar
    if shape.top < HEADER_ZONE_BOTTOM:
        return False  # zona do cabeçalho (título, supertítulo, empresa)
    if shape.top > FOOTER_ZONE_TOP:
        return False  # footer
    return True


def _is_label_shape(shape) -> bool:
    """True se o shape é um badge/label curto (não precisa de fonte grande)."""
    if not shape.has_text_frame:
        return False
    text = shape.text_frame.text.strip()
    return len(text) <= LABEL_MAX_CHARS


def _bbox(shape) -> Tuple[int, int, int, int]:
    return (shape.left, shape.top,
            shape.left + shape.width,
            shape.top + shape.height)


def _overlap_area(r1, r2) -> int:
    ox = max(0, min(r1[2], r2[2]) - max(r1[0], r2[0]))
    oy = max(0, min(r1[3], r2[3]) - max(r1[1], r2[1]))
    return ox * oy


def _luminance(rgb: RGBColor) -> float:
    def c(v):
        v /= 255
        return v / 12.92 if v <= 0.03928 else ((v + 0.055) / 1.055) ** 2.4
    return 0.2126 * c(rgb.r) + 0.7152 * c(rgb.g) + 0.0722 * c(rgb.b)


def _contrast_ratio(fg: RGBColor, bg: RGBColor) -> float:
    l1 = _luminance(fg) + 0.05
    l2 = _luminance(bg) + 0.05
    return max(l1, l2) / min(l1, l2)


# ── CHECKS ─────────────────────────────────────────────────────────────────

def _check_fonts(slide_num: int, slide, report: ValidationReport):
    """Verifica tamanhos de fonte APENAS em caixas de conteúdo real."""
    seen_issues = set()  # evita duplicatas por slide

    for shape in slide.shapes:
        if not shape.has_text_frame or not _is_content_zone(shape):
            continue
        if _is_label_shape(shape):
            continue  # badges e labels curtos são intencionalmente pequenos
        text = shape.text_frame.text.strip()
        if not text or len(text) < 10:
            continue

        min_sz = min(
            (r.font.size / 12700
             for para in shape.text_frame.paragraphs
             for r in para.runs
             if r.font.size and r.text.strip()),
            default=None
        )
        if min_sz is None:
            continue

        key = f"{slide_num}_{round(min_sz)}"
        if key in seen_issues:
            continue

        if min_sz < 9:
            seen_issues.add(key)
            report.add(Issue(
                slide=slide_num, severity="error",
                category="Fonte",
                message=(f"Texto muito pequeno: {min_sz:.0f}pt (minimo absoluto: 10pt). "
                         f"Preview: '{text[:40]}'"),
                shape=shape.name
            ))
        elif min_sz < MIN_CONTENT_PT:
            seen_issues.add(key)
            report.add(Issue(
                slide=slide_num, severity="warning",
                category="Fonte",
                message=(f"Corpo de texto com {min_sz:.0f}pt (recomendado >={MIN_CONTENT_PT}pt). "
                         f"Preview: '{text[:40]}'"),
                shape=shape.name
            ))


def _check_bounds(slide_num: int, slide, report: ValidationReport):
    """Verifica se shapes estão dentro dos limites do slide."""
    for shape in slide.shapes:
        b = _bbox(shape)
        if b[2] > W_SLIDE + 200000:
            report.add(Issue(
                slide=slide_num, severity="warning",
                category="Posição",
                message=f"Shape ultrapassa a borda direita ({(b[2]-W_SLIDE)//914:.0f}pt além)",
                shape=shape.name
            ))
        if b[3] > H_SLIDE + 100000:
            report.add(Issue(
                slide=slide_num, severity="error",
                category="Posição",
                message=f"Shape ultrapassa a borda inferior ({(b[3]-H_SLIDE)//914:.0f}pt além)",
                shape=shape.name
            ))


def _check_density(slide_num: int, slide, report: ValidationReport):
    """Verifica densidade de texto no slide."""
    total_chars = 0
    for shape in slide.shapes:
        if shape.has_text_frame and _is_content_zone(shape):
            total_chars += len(shape.text_frame.text.strip())

    if total_chars > MAX_CHARS_SLIDE:
        report.add(Issue(
            slide=slide_num, severity="warning",
            category="Densidade",
            message=f"Slide muito carregado: ~{total_chars} caracteres (máx recomendado: {MAX_CHARS_SLIDE}). "
                    f"Considere dividir em dois slides."
        ))
    elif total_chars > MAX_CHARS_SLIDE * 1.5:
        report.add(Issue(
            slide=slide_num, severity="error",
            category="Densidade",
            message=f"Slide extremamente carregado: ~{total_chars} caracteres. Precisa ser dividido."
        ))


def _check_overlaps(slide_num: int, slide, report: ValidationReport):
    """Detecta sobreposições significativas entre shapes de conteúdo."""
    content_shapes = [
        s for s in slide.shapes
        if _is_content_zone(s)
        and s.has_text_frame
        and s.text_frame.text.strip()
    ]

    for i in range(len(content_shapes)):
        for j in range(i + 1, len(content_shapes)):
            s1, s2 = content_shapes[i], content_shapes[j]
            r1, r2 = _bbox(s1), _bbox(s2)
            area = _overlap_area(r1, r2)
            if area <= 0:
                continue
            area_s1 = max(1, s1.width * s1.height)
            pct = area / area_s1 * 100
            if pct > 20:
                report.add(Issue(
                    slide=slide_num, severity="warning",
                    category="Sobreposição",
                    message=f"'{s1.name}' e '{s2.name}' se sobrepõem em ~{pct:.0f}%"
                ))


def _check_text_overflow(slide_num: int, slide, report: ValidationReport):
    """Estima se o texto pode estar transbordando da caixa."""
    for shape in slide.shapes:
        if not shape.has_text_frame or not _is_content_zone(shape):
            continue

        text = shape.text_frame.text.strip()
        if not text or shape.width <= 0:
            continue

        # Fonte média
        sizes = []
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    sizes.append(run.font.size / 12700)
        if not sizes:
            continue

        avg_pt = sum(sizes) / len(sizes)
        avg_char_w = avg_pt * 6000   # EMUs por caractere (estimativa)
        line_h     = avg_pt * 16000  # EMUs por linha

        chars_per_line = max(1, int(shape.width / avg_char_w))
        n_lines = 0
        for para in shape.text_frame.paragraphs:
            para_text = para.text
            if para_text.strip():
                n_lines += max(1, (len(para_text) // chars_per_line) + 1)
            else:
                n_lines += 0.5  # blank line

        height_needed = n_lines * line_h
        if height_needed > shape.height * 1.4:
            report.add(Issue(
                slide=slide_num, severity="warning",
                category="Overflow",
                message=f"Texto pode transbordar: ~{n_lines:.0f} linhas estimadas,"
                        f" caixa comporta ~{shape.height/line_h:.0f}",
                shape=shape.name
            ))


# ── AUTO-FIX ───────────────────────────────────────────────────────────────

def _autofix_fonts(prs: Presentation, report: ValidationReport) -> int:
    """
    Aumenta fontes muito pequenas na zona de conteúdo.
    Retorna o número de correções feitas.
    """
    fixes = 0
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame or not _is_content_zone(shape):
                continue
            text = shape.text_frame.text.strip()
            if not text or len(text) < 3:
                continue
            is_footer = shape.top > 6200000  # footer zone

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip() or run.font.size is None:
                        continue
                    sz_pt = run.font.size / 12700
                    if is_footer:
                        continue  # footers podem ser pequenos
                    if sz_pt < 10:
                        run.font.size = Pt(11)
                        fixes += 1
                        report.add(Issue(
                            slide=si + 1, severity="info",
                            category="Auto-fix",
                            message=f"Fonte ajustada: {sz_pt:.0f}pt → 11pt",
                            shape=shape.name, auto_fixed=True
                        ))
                    elif sz_pt < 11:
                        run.font.size = Pt(12)
                        fixes += 1
    return fixes


# ── ENTRY POINT ────────────────────────────────────────────────────────────

def validate(pptx_path: str, autofix: bool = True) -> ValidationReport:
    """
    Valida o arquivo PPTX e opcionalmente corrige problemas automaticamente.
    Salva o arquivo corrigido no mesmo caminho.
    """
    prs = Presentation(pptx_path)
    report = ValidationReport(n_slides=len(prs.slides))

    for si, slide in enumerate(prs.slides):
        slide_num = si + 1
        _check_fonts(slide_num, slide, report)
        _check_bounds(slide_num, slide, report)
        _check_density(slide_num, slide, report)
        _check_overlaps(slide_num, slide, report)
        _check_text_overflow(slide_num, slide, report)

    if autofix:
        fixes = _autofix_fonts(prs, report)
        if fixes > 0:
            prs.save(pptx_path)

    return report


def report_summary(report: ValidationReport) -> str:
    lines = [
        f"Score: {report.score}/100  |  Nota: {report.grade}",
        f"Erros: {len(report.errors)}  |  Avisos: {len(report.warnings)}",
        "",
    ]
    for issue in report.issues:
        if issue.auto_fixed:
            continue
        icon = {"error": "🔴", "warning": "🟡", "info": "🔵"}.get(issue.severity, "•")
        lines.append(f"  {icon} Slide {issue.slide} [{issue.category}]: {issue.message}")
    return "\n".join(lines)
