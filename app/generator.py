"""
generator.py – Converte um dicionário de pauta em .pptx com identidade visual Matriz Educação.
Suporta 3 temas: "escuro" (padrão), "claro", "premium".
"""
from __future__ import annotations
import io
import re
from dataclasses import dataclass
from pathlib import Path
from typing import List

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── PALETA BASE ─────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0E, 0x3D, 0x52)
TEAL      = RGBColor(0x5D, 0xD4, 0xB4)
DARK_NAVY = RGBColor(0x0A, 0x2D, 0x3D)
LIGHT_BG  = RGBColor(0xF4, 0xF7, 0xF8)
AMBER     = RGBColor(0xD9, 0x77, 0x06)
GREEN     = RGBColor(0x16, 0xA3, 0x4A)
RED_SOFT  = RGBColor(0xDC, 0x26, 0x26)
GRAY_TXT  = RGBColor(0x6B, 0x72, 0x80)
DARK_TXT  = RGBColor(0x1F, 0x2A, 0x37)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
SEPARATOR = RGBColor(0xE5, 0xE7, 0xEB)

# ── PARÂMETROS DE QUALIDADE (sobrescritos por quality_loop) ─────────────────
_GEN_PARAMS: dict = {}

def _get_param(key: str, default):
    return _GEN_PARAMS.get(key, default)

# ── DIMENSÕES ───────────────────────────────────────────────────────────────
W, H        = 12192000, 6858000
SIDEBAR_W   = 164592
HEADER_H    = 1097280
CONTENT_L   = 548640
CONTENT_W   = 10058400
FULL_W      = W - CONTENT_L - 200000
LOGO_L, LOGO_T = 465055, -116672
LOGO_W, LOGO_H = 2864869, 1053344
CO_L, CO_T  = 8686800, 320040
SEC_T       = 566928
SUP_T       = 1051560
TITLE_T     = 1325880
TITLE_H     = 650000
CONTENT_TOP = 2150000
FOOTER_T    = 6537960
SAFE_BOTTOM = 6450000

# ── TIPOGRAFIA ──────────────────────────────────────────────────────────────
T = {
    "title":        30,
    "card_title":   16,
    "card_body":    13,
    "bullet":       13,
    "label":        10,
    "badge":        10,
    "supertitle":   11,
    "company":      10,
    "section_lbl":   9,
    "footer":       10,
    "kpi_value":    24,
    "kpi_label":    10,
    "table_header": 11,
    "table_body":   11,
    "tag":          10,
}

# ── TEMAS ───────────────────────────────────────────────────────────────────

@dataclass
class Theme:
    name: str
    slug: str
    card_style: str             # "top_bar" | "left_bar"
    slide_bg: RGBColor | None   # None = branco implícito
    sidebar: RGBColor
    sidebar_accent: RGBColor
    title_col: RGBColor
    supertitle_col: RGBColor
    company_col: RGBColor
    separator_col: RGBColor
    footer_col: RGBColor
    card_bg: RGBColor
    card_title_col: RGBColor
    bullet_col: RGBColor
    accent_colors: List[RGBColor]
    block_bg: RGBColor
    block_accent_col: RGBColor
    block_label_col: RGBColor
    block_title_col: RGBColor
    tag_bg: RGBColor
    tag_text_col: RGBColor


THEME_ESCURO = Theme(
    name="Escuro Clássico", slug="escuro", card_style="top_bar",
    slide_bg=None,
    sidebar=NAVY, sidebar_accent=TEAL,
    title_col=DARK_NAVY, supertitle_col=TEAL,
    company_col=NAVY, separator_col=SEPARATOR, footer_col=GRAY_TXT,
    card_bg=LIGHT_BG, card_title_col=DARK_NAVY, bullet_col=DARK_TXT,
    accent_colors=[TEAL, NAVY, TEAL, NAVY, AMBER, TEAL, NAVY, AMBER],
    block_bg=DARK_NAVY, block_accent_col=TEAL,
    block_label_col=TEAL, block_title_col=WHITE,
    tag_bg=NAVY, tag_text_col=WHITE,
)

THEME_CLARO = Theme(
    name="Claro Profissional", slug="claro", card_style="left_bar",
    slide_bg=WHITE,
    sidebar=NAVY, sidebar_accent=TEAL,
    title_col=DARK_NAVY, supertitle_col=TEAL,
    company_col=DARK_NAVY, separator_col=TEAL, footer_col=GRAY_TXT,
    card_bg=WHITE, card_title_col=DARK_NAVY, bullet_col=DARK_TXT,
    accent_colors=[TEAL, NAVY, TEAL, NAVY, AMBER, TEAL, NAVY, AMBER],
    block_bg=NAVY, block_accent_col=TEAL,
    block_label_col=TEAL, block_title_col=WHITE,
    tag_bg=NAVY, tag_text_col=WHITE,
)

_D2 = RGBColor(0x07, 0x1E, 0x2B)   # fundo very dark
_D3 = RGBColor(0x05, 0x16, 0x21)   # sidebar still darker
_C3 = RGBColor(0x0D, 0x2E, 0x41)   # card bg premium
_S3 = RGBColor(0x1A, 0x45, 0x5E)   # separator premium
_F3 = RGBColor(0x7B, 0xA8, 0xBE)   # footer text premium
_B3 = RGBColor(0xC8, 0xE6, 0xDC)   # bullet text premium
_K3 = RGBColor(0x0A, 0x28, 0x3A)   # block bg premium

THEME_PREMIUM = Theme(
    name="Premium Escuro", slug="premium", card_style="top_bar",
    slide_bg=_D2,
    sidebar=_D3, sidebar_accent=TEAL,
    title_col=WHITE, supertitle_col=TEAL,
    company_col=TEAL, separator_col=_S3, footer_col=_F3,
    card_bg=_C3, card_title_col=TEAL, bullet_col=_B3,
    accent_colors=[TEAL, AMBER, TEAL, AMBER, TEAL, AMBER, TEAL, AMBER],
    block_bg=_K3, block_accent_col=AMBER,
    block_label_col=AMBER, block_title_col=WHITE,
    tag_bg=_S3, tag_text_col=WHITE,
)

THEMES = {
    "escuro":  THEME_ESCURO,
    "claro":   THEME_CLARO,
    "premium": THEME_PREMIUM,
}

# ── HELPERS PRIMITIVOS ──────────────────────────────────────────────────────

def _load_logo(assets_dir: Path) -> bytes | None:
    for name in ("logo.png", "logo.jpg", "logo.jpeg"):
        p = assets_dir / name
        if p.exists():
            return p.read_bytes()
    return None


def _rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Emu(l), Emu(t), Emu(w), Emu(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _txt(slide, l, t, w, h, text, sz,
         bold=False, color=DARK_TXT, align=PP_ALIGN.LEFT,
         italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text        = text
    r.font.size   = Pt(sz)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name   = "Calibri"
    return tb


def _bullets(slide, l, t, w, h, items,
             sz=None, color=DARK_TXT, spacing=8, max_items=None):
    if not items:
        return
    sz = sz or T["bullet"]
    display = list(items)
    if max_items and len(display) > max_items:
        display = display[:max_items]
        display.append(f"… e mais {len(items) - max_items} itens")

    tb = slide.shapes.add_textbox(Emu(l), Emu(t), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(display):
        if isinstance(item, dict):
            text = item.get("text", "")
            s    = item.get("sz", sz)
            c    = item.get("color", color)
            bd   = item.get("bold", False)
            bul  = item.get("bullet", True)
        else:
            text = str(item)
            s, c, bd, bul = sz, color, False, True
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(spacing)
        r = p.add_run()
        r.text       = ("• " if bul else "") + text
        r.font.size  = Pt(s)
        r.font.bold  = bd
        r.font.color.rgb = c
        r.font.name  = "Calibri"


def _add_logo(slide, logo_blob):
    if logo_blob:
        slide.shapes.add_picture(
            io.BytesIO(logo_blob),
            Emu(LOGO_L), Emu(LOGO_T), Emu(LOGO_W), Emu(LOGO_H))


def _tag(slide, l, t, w, h, text, bg, text_col=WHITE):
    _rect(slide, l, t, w, h, bg)
    _txt(slide, l, t, w, h, text, T["tag"], bold=True, color=text_col,
         align=PP_ALIGN.CENTER)


def _card(slide, l, t, w, h, accent, theme: Theme):
    """Card com fundo e acento conforme o estilo do tema."""
    if theme.card_style == "left_bar":
        _rect(slide, l, t, w, h, theme.card_bg)
        _rect(slide, l, t, 80000, h, accent)           # barra vertical esquerda
        _rect(slide, l, t, w, 18288, theme.separator_col)  # linha fina no topo
    else:
        _rect(slide, l, t, w, h, theme.card_bg)
        _rect(slide, l, t, w, 73152, accent)            # barra horizontal no topo


def _header(slide, logo_blob, supertitle: str, title: str,
            theme: Theme, section: str = "Reunião de Diretores",
            footer: str = ""):
    if theme.slide_bg:
        _rect(slide, 0, 0, W, H, theme.slide_bg)
    _rect(slide, 0, 0, SIDEBAR_W, H, theme.sidebar)
    _rect(slide, 0, 0, SIDEBAR_W, HEADER_H, theme.sidebar_accent)
    _txt(slide, CO_L, CO_T, 3200400, 274320,
         "MATRIZ EDUCAÇÃO", T["company"], bold=True,
         color=theme.company_col, align=PP_ALIGN.RIGHT)
    _txt(slide, CO_L, SEC_T, 3200400, 228600,
         section, T["section_lbl"], color=theme.footer_col, align=PP_ALIGN.RIGHT)
    _rect(slide, CONTENT_L, 868680, 11338560, 9144, theme.separator_col)
    _txt(slide, CONTENT_L, SUP_T, FULL_W, 274320,
         supertitle.upper(), T["supertitle"], bold=True, color=theme.supertitle_col)
    _txt(slide, CONTENT_L, TITLE_T, FULL_W, TITLE_H,
         title, T["title"], bold=True, color=theme.title_col)
    if footer:
        _txt(slide, CONTENT_L, FOOTER_T, 7315200, 274320,
             footer, T["footer"], color=theme.footer_col)
    _add_logo(slide, logo_blob)


# ══════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════

def _slide_capa(prs, blank, logo_blob, pauta: dict, footer: str, theme: Theme):
    s = prs.slides.add_slide(blank)

    if theme.slug == "claro":
        _rect(s, 0, 0, W, H, WHITE)
        _rect(s, 0, 0, SIDEBAR_W, H, NAVY)
        _rect(s, 0, 0, SIDEBAR_W, HEADER_H, TEAL)
        _rect(s, W - 1500000, 0, 1500000, H, LIGHT_BG)
        _rect(s, W - 400000,  0, 400000,  H, TEAL)
        _rect(s, W - 800000,  0, 9144,    H, SEPARATOR)
        _rect(s, CONTENT_L, 1500000, W - SIDEBAR_W - 1800000, 1000000, NAVY)
        _rect(s, CONTENT_L, 1500000, 120000, 1000000, TEAL)
        capa_card_bg  = LIGHT_BG
        capa_border   = TEAL
        capa_lbl_col  = TEAL
        capa_val_col  = DARK_NAVY
        agenda_txt    = DARK_TXT
    elif theme.slug == "premium":
        _rect(s, 0, 0, W, H, _D2)
        _rect(s, SIDEBAR_W, 0, W - SIDEBAR_W, H, RGBColor(0x09, 0x25, 0x36))
        _rect(s, 0, 0, SIDEBAR_W, H, _D3)
        _rect(s, 0, 0, SIDEBAR_W, HEADER_H, TEAL)
        _rect(s, W - 1500000, 0, 1500000, H, RGBColor(0x05, 0x1A, 0x28))
        _rect(s, W - 400000,  0, 400000,  H, TEAL)
        _rect(s, W - 800000,  0, 9144,    H, _S3)
        _rect(s, CONTENT_L, 1500000, W - SIDEBAR_W - 1800000, 1000000, _K3)
        _rect(s, CONTENT_L, 1500000, 120000, 1000000, TEAL)
        capa_card_bg  = _C3
        capa_border   = AMBER
        capa_lbl_col  = AMBER
        capa_val_col  = WHITE
        agenda_txt    = WHITE
    else:  # escuro
        _rect(s, 0, 0, W, H, NAVY)
        _rect(s, SIDEBAR_W, 0, W - SIDEBAR_W, H, WHITE)
        _rect(s, 0, 0, SIDEBAR_W, H, NAVY)
        _rect(s, 0, 0, SIDEBAR_W, HEADER_H, TEAL)
        _rect(s, W - 1500000, 0, 1500000, H, RGBColor(0xF0, 0xFB, 0xF7))
        _rect(s, W - 400000,  0, 400000,  H, TEAL)
        _rect(s, W - 800000,  0, 9144,    H, SEPARATOR)
        _rect(s, CONTENT_L, 1500000, W - SIDEBAR_W - 1800000, 1000000, DARK_NAVY)
        _rect(s, CONTENT_L, 1500000, 120000, 1000000, TEAL)
        capa_card_bg  = LIGHT_BG
        capa_border   = TEAL
        capa_lbl_col  = TEAL
        capa_val_col  = DARK_NAVY
        agenda_txt    = DARK_TXT

    _add_logo(s, logo_blob)
    _txt(s, CO_L, CO_T, 3000000, 274320,
         "MATRIZ EDUCAÇÃO", T["company"], bold=True,
         color=theme.company_col, align=PP_ALIGN.RIGHT)
    _rect(s, CONTENT_L, 868680, W - SIDEBAR_W - 1600000, 9144, theme.separator_col)

    _txt(s, CONTENT_L + 220000, 1570000, 8000000, 260000,
         "REUNIÃO DE DIRETORES", T["supertitle"], bold=True, color=TEAL)
    data_fmt = pauta.get("data", "") or "—"
    _txt(s, CONTENT_L + 220000, 1820000, 8000000, 620000,
         data_fmt, 32, bold=True, color=WHITE)

    if pauta.get("objetivo"):
        _txt(s, CONTENT_L, 2650000, W - SIDEBAR_W - 1800000, 350000,
             pauta["objetivo"], 13, color=GRAY_TXT, italic=True, wrap=True)

    _rect(s, CONTENT_L, 3100000, W - SIDEBAR_W - 1800000, 9144, theme.separator_col)

    agenda = pauta.get("agenda", [])
    total_min = sum(
        int(re.search(r"\d+", a["tempo"]).group())
        for a in agenda if a.get("tempo") and re.search(r"\d+", a["tempo"])
    )
    for i, (lbl, val) in enumerate([
        ("DURAÇÃO TOTAL", f"{total_min} min" if total_min else "—"),
        ("FORMATO",       "Presencial · Reunião de Diretores"),
    ]):
        ox = CONTENT_L + i * 5100000
        _rect(s, ox, 3250000, 4800000, 560000, capa_card_bg)
        _rect(s, ox, 3250000, 80000, 560000, capa_border)
        _txt(s, ox + 180000, 3275000, 4400000, 240000, lbl, T["label"], bold=True, color=capa_lbl_col)
        _txt(s, ox + 180000, 3530000, 4400000, 280000, val, 13, color=capa_val_col)

    _rect(s, CONTENT_L, 3950000, W - SIDEBAR_W - 1800000, 9144, TEAL)
    _txt(s, CONTENT_L, 4010000, 5000000, 240000, "PAUTA", T["supertitle"], bold=True, color=TEAL)

    real_items = [a for a in agenda if a["num"] > 0]
    half = (len(real_items) + 1) // 2
    for ci, chunk in enumerate([real_items[:half], real_items[half:]]):
        lx = CONTENT_L + ci * 5100000
        for ri, item in enumerate(chunk):
            _txt(s, lx, 4280000 + ri * 390000, 4800000, 380000,
                 f"{item['num']:02d}  {item['nome']}" +
                 (f"  —  {item['tempo']}" if item.get("tempo") else ""),
                 13, bold=True, color=agenda_txt)

    _txt(s, CONTENT_L, FOOTER_T, 7315200, 274320, footer, T["footer"], color=theme.footer_col)


def _slide_agenda(prs, blank, logo_blob, pauta: dict, footer: str, theme: Theme):
    s = prs.slides.add_slide(blank)
    _header(s, logo_blob,
            f"Reunião de Diretores · {pauta.get('data', '')}",
            "Pauta do Dia", theme=theme, footer=footer)

    agenda  = pauta.get("agenda", [])
    ROW_H   = 530000
    ROW_GAP = 60000
    y = CONTENT_TOP

    for item in agenda:
        if item["num"] == 0:
            _rect(s, CONTENT_L, y + ROW_GAP // 2, FULL_W, 280000, theme.card_bg)
            _rect(s, CONTENT_L, y + ROW_GAP // 2, 80000, 280000, theme.separator_col)
            _txt(s, CONTENT_L + 170000, y + ROW_GAP // 2 + 50000,
                 5000000, 220000,
                 f"  INTERVALO  —  {item['tempo']}", T["label"],
                 bold=True, color=theme.footer_col)
            y += 280000 + ROW_GAP
            continue

        accent = theme.accent_colors[(item["num"] - 1) % len(theme.accent_colors)]
        _rect(s, CONTENT_L, y, FULL_W, ROW_H, theme.card_bg)
        _rect(s, CONTENT_L, y, 80000, ROW_H, accent)
        _txt(s, CONTENT_L + 80000, y + 10000, 380000, ROW_H - 20000,
             f"{item['num']:02d}", 16, bold=True, color=accent,
             align=PP_ALIGN.CENTER)
        _txt(s, CONTENT_L + 520000, y + 90000, 8000000, 280000,
             item["nome"], 16, bold=True, color=theme.title_col)
        if item.get("tempo"):
            _tag(s,
                 CONTENT_L + FULL_W - 1200000,
                 y + (ROW_H - 310000) // 2,
                 1150000, 310000,
                 item["tempo"], theme.tag_bg, theme.tag_text_col)
        y += ROW_H + ROW_GAP

    _rect(s, CONTENT_L, SAFE_BOTTOM - 310000, FULL_W, 300000, TEAL)
    _txt(s, CONTENT_L + 200000, SAFE_BOTTOM - 200000, 5000000, 280000,
         "ENCERRAMENTO  ·  FIM DA REUNIÃO", T["label"],
         bold=True, color=DARK_NAVY)


def _slide_with_subs(prs, blank, logo_blob, sec, supertitle, accent,
                     subsecoes, extra_items, footer, theme: Theme):
    s = prs.slides.add_slide(blank)
    _header(s, logo_blob, supertitle, sec["titulo"], theme=theme, footer=footer)

    n    = min(len(subsecoes), 4)
    cols = 2 if n > 2 else n
    rows = (n + 1) // 2

    CARD_GAP   = 80000
    BW         = (FULL_W - (cols - 1) * CARD_GAP) // cols
    avail_h    = SAFE_BOTTOM - CONTENT_TOP
    BH_natural = (avail_h - (rows - 1) * CARD_GAP) // rows
    BH         = min(BH_natural, 2600000)
    BH         = max(BH, 1400000)

    for i, sub in enumerate(subsecoes[:4]):
        row, col = divmod(i, cols)
        bx = CONTENT_L + col * (BW + CARD_GAP)
        by = CONTENT_TOP + row * (BH + CARD_GAP)

        _card(s, bx, by, BW, BH, accent, theme)

        if theme.card_style == "top_bar":
            _rect(s, bx, by + 73152, 80000, BH - 73152, theme.card_bg)
            title_top = by + 130000
        else:
            title_top = by + 100000

        title_text = sub["titulo"]
        title_h    = 450000 if len(title_text) > 35 else 320000
        _txt(s, bx + 170000, title_top, BW - 270000, title_h,
             title_text, T["card_title"], bold=True,
             color=theme.card_title_col, wrap=True)

        items_disp     = sub.get("items", [])
        max_items      = _get_param('max_bullets_per_card',
                                    max(3, int((BH - 600000) // (T["bullet"] * 20000))))
        spacing_factor = _get_param('extra_spacing_factor', 1.0)
        bullets_top    = title_top + title_h + 60000
        bullets_h      = by + BH - 80000 - bullets_top
        _bullets(s, bx + 170000, bullets_top,
                 BW - 280000, max(bullets_h, 500000),
                 items_disp, sz=T["bullet"], color=theme.bullet_col,
                 spacing=int(8 * spacing_factor), max_items=max_items)

    if len(subsecoes) > 4:
        _slide_with_subs(prs, blank, logo_blob, sec, supertitle, accent,
                         subsecoes[4:], extra_items, footer, theme)


def _slide_bullets_only(prs, blank, logo_blob, sec, supertitle, accent,
                        items, footer, theme: Theme):
    MAX_PER_SLIDE = 8
    for chunk_start in range(0, len(items), MAX_PER_SLIDE):
        chunk = items[chunk_start:chunk_start + MAX_PER_SLIDE]
        s = prs.slides.add_slide(blank)
        part_label = supertitle + (" (cont.)" if chunk_start > 0 else "")
        _header(s, logo_blob, part_label, sec["titulo"], theme=theme, footer=footer)

        CT = CONTENT_TOP

        if sec.get("subtitulo") and sec["subtitulo"] != sec["titulo"] and chunk_start == 0:
            obj_txt = sec.get("objetivo_label", "") or sec.get("subtitulo", "")
            block_h = 850000 if (obj_txt and len(obj_txt) > 60) else 730000
            _rect(s, CONTENT_L, CT, FULL_W, block_h, theme.block_bg)
            _rect(s, CONTENT_L, CT, 110000, block_h, theme.block_accent_col)
            _txt(s, CONTENT_L + 210000, CT + 120000, FULL_W - 300000, 220000,
                 "OBJETIVO", T["label"], bold=True, color=theme.block_label_col)
            _txt(s, CONTENT_L + 210000, CT + 360000, FULL_W - 300000, block_h - 370000,
                 sec["subtitulo"], T["card_title"], bold=True,
                 color=theme.block_title_col, wrap=True)
            CT += block_h + 100000

        avail = SAFE_BOTTOM - CT - 100000

        if len(chunk) > 4:
            half  = (len(chunk) + 1) // 2
            col_w = (FULL_W - 150000) // 2
            for ci, col_items in enumerate([chunk[:half], chunk[half:]]):
                lx = CONTENT_L + ci * (col_w + 150000)
                _bullets(s, lx, CT + 150000, col_w, avail - 150000,
                         col_items, sz=T["bullet"], color=theme.bullet_col, spacing=10)
        else:
            _bullets(s, CONTENT_L, CT + 150000, FULL_W, avail - 150000,
                     chunk, sz=T["bullet"], color=theme.bullet_col, spacing=14)


def _slide_section_generic(prs, blank, logo_blob, sec: dict,
                           footer: str, idx: int, theme: Theme):
    accent     = theme.accent_colors[idx % len(theme.accent_colors)]
    supertitle = f"{sec['num']:02d} · {sec['titulo']}"
    if sec.get("tempo"):
        supertitle += f"  —  {sec['tempo']}"

    items_all = [it for it in sec.get("items", []) if it.strip()]
    subsecoes = [sub for sub in sec.get("subsecoes", [])
                 if sub.get("items") or sub.get("titulo")]

    if subsecoes:
        _slide_with_subs(prs, blank, logo_blob, sec, supertitle, accent,
                         subsecoes, items_all, footer, theme)
    elif items_all:
        _slide_bullets_only(prs, blank, logo_blob, sec, supertitle, accent,
                            items_all, footer, theme)
    else:
        s = prs.slides.add_slide(blank)
        _header(s, logo_blob, supertitle, sec["titulo"], theme=theme, footer=footer)
        _txt(s, CONTENT_L, CONTENT_TOP + 400000, FULL_W, 300000,
             "Conteúdo a ser apresentado na reunião",
             T["card_body"], color=theme.footer_col, italic=True)


def _slide_encerramento(prs, blank, logo_blob, pauta: dict, footer: str, theme: Theme):
    s = prs.slides.add_slide(blank)

    enc_bg = _D2 if theme.slug == "premium" else NAVY
    _rect(s, 0, 0, W, H, enc_bg)
    _rect(s, 0, 0, SIDEBAR_W, H, theme.sidebar)
    _rect(s, 0, 0, SIDEBAR_W, HEADER_H, theme.sidebar_accent)
    _rect(s, W - 400000, 0, 400000, H, TEAL)
    _add_logo(s, logo_blob)

    _txt(s, CO_L, CO_T, 3200400, 274320,
         "MATRIZ EDUCAÇÃO", T["company"], bold=True,
         color=TEAL, align=PP_ALIGN.RIGHT)

    block_col = _K3 if theme.slug == "premium" else RGBColor(0x16, 0x4E, 0x63)
    _rect(s, CONTENT_L, 1900000, FULL_W, 1000000, block_col)
    _rect(s, CONTENT_L, 1900000, 120000, 1000000, TEAL)
    _txt(s, CONTENT_L + 220000, 1970000, FULL_W - 300000, 260000,
         "ENCERRAMENTO", T["supertitle"], bold=True, color=TEAL)
    _txt(s, CONTENT_L + 220000, 2220000, FULL_W - 300000, 580000,
         f"Reunião de Diretores  —  {pauta.get('data', '')}",
         24, bold=True, color=WHITE)

    enc_items = [
        ("Dúvidas",                "Espaço aberto para perguntas e esclarecimentos"),
        ("Temas Diversos",         "Assuntos não contemplados na pauta"),
        ("Reforço dos Combinados", "Revisão dos acordos e responsabilidades da reunião"),
        ("Prioridades da Semana",  "Alinhamento do que cada área executará nos próximos dias"),
    ]
    EW  = (FULL_W - 80000) // 2
    EH  = 780000
    enc_card = _C3 if theme.slug == "premium" else RGBColor(0x16, 0x4E, 0x63)
    for i, (title, desc) in enumerate(enc_items):
        row, col = divmod(i, 2)
        ex = CONTENT_L + col * (EW + 80000)
        ey = 3100000 + row * (EH + 73152)
        acc = theme.accent_colors[i % len(theme.accent_colors)]
        _rect(s, ex, ey, EW, EH, enc_card)
        _rect(s, ex, ey, EW, 73152, acc)
        _rect(s, ex, ey + 73152, 80000, EH - 73152, enc_bg)
        _txt(s, ex + 170000, ey + 140000, EW - 250000, 300000,
             title, T["card_title"], bold=True, color=TEAL, wrap=True)
        _txt(s, ex + 170000, ey + 460000, EW - 250000, 300000,
             desc, T["card_body"], color=_B3, wrap=True)

    _txt(s, CONTENT_L, FOOTER_T, 7315200, 274320, footer, T["footer"], color=theme.footer_col)
    _txt(s, 8686800, FOOTER_T, 3200400, 274320,
         "FIM", T["footer"], bold=True, color=TEAL, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════
# ENTRY POINT
# ══════════════════════════════════════════════════════════════════════════

def generate(pauta: dict, output_path: str,
             assets_dir: str | Path | None = None,
             theme_name: str = "escuro") -> str:
    global _GEN_PARAMS
    _GEN_PARAMS = pauta.get('_gen_params', {})

    if assets_dir is None:
        assets_dir = Path(__file__).parent / "assets"
    assets_dir = Path(assets_dir)
    logo_blob  = _load_logo(assets_dir)

    theme = THEMES.get(theme_name, THEME_ESCURO)

    prs = Presentation()
    prs.slide_width  = Emu(W)
    prs.slide_height = Emu(H)
    blank = prs.slide_layouts[6]

    data   = pauta.get("data", "")
    footer = f"Matriz Educação  ·  Reunião de Diretores{' — ' + data if data else ''}"

    _slide_capa(prs, blank, logo_blob, pauta, footer, theme)

    if pauta.get("agenda"):
        _slide_agenda(prs, blank, logo_blob, pauta, footer, theme)

    for i, sec in enumerate(pauta.get("secoes", [])):
        _slide_section_generic(prs, blank, logo_blob, sec, footer, i, theme)

    _slide_encerramento(prs, blank, logo_blob, pauta, footer, theme)

    prs.save(output_path)
    return output_path
