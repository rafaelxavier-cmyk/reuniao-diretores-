"""
layout_audit.py – Auditor de layout de precisão.
Detecta overlaps, text overflow e elementos fora de posição.
"""
from __future__ import annotations
from dataclasses import dataclass, field
from typing import List
from pptx import Presentation
from pptx.util import Emu

SLIDE_W   = 12192000
SLIDE_H   = 6858000
SIDEBAR_W = 164592
FOOTER_TOP = 6430000   # zona de footer (abaixo deste Y é esperado ter elementos pequenos)
CONTENT_L  = 548640
CONTENT_TOP = 900000   # tudo acima disso é cabeçalho estrutural

@dataclass
class LayoutIssue:
    slide: int
    kind: str          # 'overflow' | 'overlap' | 'out_of_bounds' | 'tiny_box'
    severity: str      # 'error' | 'warning'
    shape_name: str
    details: str
    auto_fixable: bool = False
    fix_hint: dict = field(default_factory=dict)  # {attr: new_value}

@dataclass
class AuditReport:
    issues: List[LayoutIssue] = field(default_factory=list)

    @property
    def score(self) -> int:
        errors   = sum(1 for i in self.issues if i.severity == 'error')
        warnings = sum(1 for i in self.issues if i.severity == 'warning')
        return max(0, 100 - errors * 15 - warnings * 5)

    @property
    def passed(self) -> bool:
        return self.score >= 85 and not any(i.severity == 'error' for i in self.issues)

    def summary(self) -> str:
        errors   = [i for i in self.issues if i.severity == 'error']
        warnings = [i for i in self.issues if i.severity == 'warning']
        lines = [f"Score: {self.score}/100  |  Erros: {len(errors)}  Avisos: {len(warnings)}"]
        for iss in self.issues:
            icon = '[ERROR]' if iss.severity == 'error' else '[WARN]'
            lines.append(f"  {icon} Slide {iss.slide} [{iss.kind}]: {iss.details}")
        return '\n'.join(lines)

def _is_structural(shape) -> bool:
    """True para shapes que são decoração estrutural (sidebar, logo, background)."""
    if shape.left < SIDEBAR_W + 50000:
        return True   # sidebar
    if shape.top < 50000 and shape.width > SLIDE_W * 0.3:
        return True   # background full-width
    if shape.top < 200000 and shape.height > 800000:
        return True   # logo (tall, near top)
    return False

def _is_footer_zone(shape) -> bool:
    return shape.top >= FOOTER_TOP

def _estimate_text_height(shape, default_pt: float = 12.0) -> int:
    """Estimate EMU height needed to display the text in this shape."""
    if not shape.has_text_frame:
        return 0
    if shape.width <= 0:
        return 0

    max_pt = default_pt
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                pt = run.font.size / 12700
                if pt > max_pt:
                    max_pt = pt

    avg_char_w = max_pt * 6200  # EMU per char (calibri, normal)
    chars_per_line = max(1, int(shape.width / avg_char_w))

    total_lines = 0
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if text:
            total_lines += max(1, (len(text) // chars_per_line) + 1)
        else:
            total_lines += 0.4  # empty paragraph adds spacing

    line_height = max_pt * 16000  # EMU per line with spacing
    return int(total_lines * line_height) + int(max_pt * 8000)  # bottom padding

def audit(pptx_path: str) -> AuditReport:
    prs = Presentation(pptx_path)
    report = AuditReport()

    for si, slide in enumerate(prs.slides):
        sn = si + 1
        all_shapes = list(slide.shapes)

        # Filter to content shapes (text boxes in content zone, not structural)
        content_txt_shapes = [
            s for s in all_shapes
            if s.has_text_frame
            and not _is_structural(s)
            and not _is_footer_zone(s)
            and s.text_frame.text.strip()
            and s.top > CONTENT_TOP
        ]

        for sh in content_txt_shapes:
            txt = sh.text_frame.text.strip()
            if not txt:
                continue

            # 1. Tiny box check: height < font size
            min_font_pt = min(
                (r.font.size / 12700 for p in sh.text_frame.paragraphs
                 for r in p.runs if r.font.size),
                default=12.0
            )
            min_viable_h = int(min_font_pt * 14000)  # at least 1 line
            if sh.height < min_viable_h:
                report.issues.append(LayoutIssue(
                    slide=sn, kind='tiny_box', severity='error',
                    shape_name=sh.name,
                    details=f"Caixa {sh.height//914:.0f}pt mas fonte {min_font_pt:.0f}pt — texto invisível. Preview: '{txt[:40]}'",
                    auto_fixable=True,
                    fix_hint={'min_height': int(min_font_pt * 20000)}
                ))

            # 2. Text overflow check
            needed_h = _estimate_text_height(sh, min_font_pt)
            if needed_h > sh.height * 1.4 and sh.height > 0:
                report.issues.append(LayoutIssue(
                    slide=sn, kind='overflow', severity='warning',
                    shape_name=sh.name,
                    details=(f"Overflow estimado: caixa {sh.height//914:.0f}pt, "
                             f"texto precisa ~{needed_h//914:.0f}pt. Preview: '{txt[:40]}'"),
                    auto_fixable=True,
                    fix_hint={'min_height': needed_h}
                ))

            # 3. Out of bounds (real content, not footer)
            bottom = sh.top + sh.height
            if bottom > SLIDE_H - 50000:
                report.issues.append(LayoutIssue(
                    slide=sn, kind='out_of_bounds', severity='error',
                    shape_name=sh.name,
                    details=f"Shape ultrapassa o slide: bottom={bottom//914:.0f}pt > {SLIDE_H//914:.0f}pt",
                ))

        # 4. Overlap detection among content text shapes
        for i in range(len(content_txt_shapes)):
            for j in range(i + 1, len(content_txt_shapes)):
                s1, s2 = content_txt_shapes[i], content_txt_shapes[j]
                l1, t1 = s1.left, s1.top
                r1, b1 = l1 + s1.width, t1 + s1.height
                l2, t2 = s2.left, s2.top
                r2, b2 = l2 + s2.width, t2 + s2.height

                ox = max(0, min(r1, r2) - max(l1, l2))
                oy = max(0, min(b1, b2) - max(t1, t2))
                area = ox * oy
                if area <= 0:
                    continue

                min_area = min(s1.width * s1.height, s2.width * s2.height)
                pct = area / max(1, min_area) * 100

                if pct > 30:
                    t1p = s1.text_frame.text.strip()[:25]
                    t2p = s2.text_frame.text.strip()[:25]
                    report.issues.append(LayoutIssue(
                        slide=sn, kind='overlap', severity='error',
                        shape_name=f"{s1.name} <-> {s2.name}",
                        details=f"Sobreposicao {pct:.0f}%: '{t1p}' (t={t1//914:.0f} b={b1//914:.0f}) e '{t2p}' (t={t2//914:.0f} b={b2//914:.0f})"
                    ))
                elif pct > 15:
                    t1p = s1.text_frame.text.strip()[:25]
                    t2p = s2.text_frame.text.strip()[:25]
                    report.issues.append(LayoutIssue(
                        slide=sn, kind='overlap', severity='warning',
                        shape_name=f"{s1.name} <-> {s2.name}",
                        details=f"Sobreposicao leve {pct:.0f}%: '{t1p}' e '{t2p}'"
                    ))

    return report
